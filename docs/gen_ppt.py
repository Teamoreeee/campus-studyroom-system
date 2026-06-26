#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《校园自习室预约系统》答辩 PPT —— 简约干净 + 蓝色风格。
依赖：python-pptx；用法：python gen_ppt.py
输出：docs/答辩PPT.pptx（16:9，24 页）
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "images")
OUT = os.path.join(BASE, "答辩PPT.pptx")

# ---------------- 配色 ----------------
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # 主蓝
DARKBLUE  = RGBColor(0x1E, 0x3A, 0x8A)   # 深蓝
MIDBLUE   = RGBColor(0x3B, 0x82, 0xF6)   # 中蓝
LIGHTBLUE = RGBColor(0xDB, 0xEA, 0xFE)   # 浅蓝
PALEBLUE  = RGBColor(0xEF, 0xF6, 0xFF)   # 极浅蓝（表格隔行/占位）
INK       = RGBColor(0x1F, 0x29, 0x37)   # 正文深灰
SUBINK    = RGBColor(0x6B, 0x72, 0x80)   # 次要灰
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "微软雅黑"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_page = 0


# ---------------- 基础工具 ----------------
def _set_font(run, size=None, bold=None, color=None, name=FONT):
    f = run.font
    f.name = name
    # 中文字体
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {})
        rpr.append(ea)
    ea.set('typeface', name)
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, l, t, w, h, color, line=None, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def round_rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_para(p, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
             space_after=6, name=FONT):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    _set_font(p.runs[0], size, bold, color, name)
    return p


# ---------------- 通用装饰 ----------------
def decorate(slide, light=False):
    """每页统一装饰：左侧细蓝条 + 右上角几何块 + 底部细线。"""
    # 左侧竖条
    rect(slide, 0, 0, Inches(0.13), SH, BLUE)
    # 右上角几何装饰
    rect(slide, SW - Inches(1.7), 0, Inches(1.7), Inches(0.28), DARKBLUE)
    rect(slide, SW - Inches(0.9), Inches(0.28), Inches(0.9), Inches(0.16), MIDBLUE)
    rect(slide, SW - Inches(2.5), 0, Inches(0.8), Inches(0.16),
         LIGHTBLUE if not light else MIDBLUE)
    # 底部细线
    rect(slide, Inches(0.6), SH - Inches(0.52), SW - Inches(1.2), Pt(1.3), LIGHTBLUE)


def footer(slide):
    global _page
    tb, tf = textbox(slide, Inches(0.6), SH - Inches(0.48), SW - Inches(1.2), Inches(0.32),
                     MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_para(p, "校园自习室预约系统  ·  软件架构技术课程设计答辩", 9, False, SUBINK)
    # 页码
    tb2, tf2 = textbox(slide, SW - Inches(2.0), SH - Inches(0.48), Inches(1.4), Inches(0.32),
                       MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_para(p2, f"{_page:02d} / 24", 9, True, BLUE)


def title(slide, text, kicker=None):
    """内容页标题：可选小标签 + 主标题 + 蓝色下划短线。"""
    top = Inches(0.55)
    if kicker:
        tb, tf = textbox(slide, Inches(0.7), Inches(0.42), Inches(8), Inches(0.3))
        set_para(tf.paragraphs[0], kicker, 12, True, MIDBLUE)
        top = Inches(0.78)
    tb, tf = textbox(slide, Inches(0.7), top, Inches(11.8), Inches(0.8))
    set_para(tf.paragraphs[0], text, 27, True, DARKBLUE)
    # 下划线
    rect(slide, Inches(0.72), top + Inches(0.78), Inches(1.5), Pt(3.5), BLUE)


def new_slide(decorated=True, footed=True):
    global _page
    s = prs.slides.add_slide(BLANK)
    _page += 1
    if decorated:
        decorate(s)
    if footed:
        footer(s)
    return s


# ---------------- 内容组件 ----------------
def bullets(slide, items, left, top, width, height, size=15, gap=10):
    """items: list of (text, level)；level=0 主点，1 子点。"""
    tb, tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        text, level = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if level == 0 else gap - 3)
        p.alignment = PP_ALIGN.LEFT
        if level == 0:
            r = p.add_run(); r.text = "▪  "
            _set_font(r, size, True, BLUE)
            r2 = p.add_run(); r2.text = text
            _set_font(r2, size, False, INK)
        else:
            r = p.add_run(); r.text = "        –  "
            _set_font(r, size - 1, False, MIDBLUE)
            r2 = p.add_run(); r2.text = text
            _set_font(r2, size - 1, False, SUBINK)
    return tb


def table(slide, headers, rows, left, top, width, height,
          hsize=13, bsize=12, col_widths=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    gt = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gt.table
    # 关闭默认样式条带，自己上色
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    # 表头
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = Pt(3); c.margin_bottom = Pt(3)
        c.margin_left = Pt(7); c.margin_right = Pt(5)
        p = c.text_frame.paragraphs[0]
        set_para(p, h, hsize, True, WHITE, PP_ALIGN.LEFT, 0)
    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else PALEBLUE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.margin_left = Pt(7); c.margin_right = Pt(5)
            p = c.text_frame.paragraphs[0]
            set_para(p, str(val), bsize, False, INK, PP_ALIGN.LEFT, 0)
    return tbl


def add_image(slide, fname, left, top, max_w, max_h, caption=None):
    path = os.path.join(IMG, fname)
    if not os.path.exists(path):
        placeholder(slide, left, top, max_w, max_h, f"[缺图] {fname}")
        return
    pic = slide.shapes.add_picture(path, left, top, width=max_w)
    if pic.height > max_h:
        factor = max_h / pic.height
        pic.width = int(pic.width * factor)
        pic.height = int(max_h)
    # 水平居中于给定区域
    pic.left = int(left + (max_w - pic.width) / 2)
    if caption:
        tb, tf = textbox(slide, left, top + pic.height + Inches(0.05), max_w, Inches(0.3))
        set_para(tf.paragraphs[0], caption, 11, False, SUBINK, PP_ALIGN.CENTER)
    return pic


def add_screenshot(slide, path, left, top, max_w, max_h, caption=None):
    if not os.path.exists(path):
        placeholder(slide, left, top, max_w, max_h, f"[缺图] {os.path.basename(path)}")
        return
    pic = slide.shapes.add_picture(path, left, top, width=max_w)
    if pic.height > max_h:
        factor = max_h / pic.height
        pic.width = int(pic.width * factor)
        pic.height = int(max_h)
    pic.left = int(left + (max_w - pic.width) / 2)
    if caption:
        tb, tf = textbox(slide, left, top + pic.height + Inches(0.05), max_w, Inches(0.3))
        set_para(tf.paragraphs[0], caption, 11, False, SUBINK, PP_ALIGN.CENTER)
    return pic


def placeholder(slide, left, top, w, h, text):
    sp = round_rect(slide, left, top, w, h, PALEBLUE, line=MIDBLUE)
    sp.line.dash_style = None
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_para(p, "🖼  " + text, 13, True, MIDBLUE, PP_ALIGN.CENTER, 0)
    return sp


def cards(slide, items, left, top, total_w, card_h, gap=Inches(0.25)):
    """items: list of (num, title, desc)；横向卡片。"""
    n = len(items)
    cw = int((total_w - gap * (n - 1)) / n)
    for i, (num, ttl, desc) in enumerate(items):
        x = int(left + i * (cw + gap))
        card = round_rect(slide, x, top, cw, card_h, WHITE, line=LIGHTBLUE)
        # 顶部色条
        rect(slide, x, top, cw, Inches(0.12), BLUE)
        # 序号圆
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), top + Inches(0.28),
                                       Inches(0.55), Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = LIGHTBLUE
        _no_line(badge); badge.shadow.inherit = False
        bp = badge.text_frame.paragraphs[0]
        set_para(bp, str(num), 18, True, DARKBLUE, PP_ALIGN.CENTER, 0)
        # 标题 + 描述
        tb, tf = textbox(slide, x + Inches(0.2), top + Inches(1.0),
                         cw - Inches(0.4), card_h - Inches(1.1))
        set_para(tf.paragraphs[0], ttl, 14, True, DARKBLUE, space_after=4)
        p = tf.add_paragraph()
        set_para(p, desc, 11, False, SUBINK)


# ==================================================================
# 以下为 24 页内容
# ==================================================================

# ---- P1 封面 ----
def slide_cover():
    global _page
    s = prs.slides.add_slide(BLANK); _page += 1
    # 背景大色块
    rect(s, 0, 0, SW, SH, WHITE)
    # 左侧深蓝竖带
    rect(s, 0, 0, Inches(0.25), SH, BLUE)
    # 右下角几何装饰
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, SW - Inches(4.2), SH - Inches(3.2),
                             Inches(4.2), Inches(3.2))
    tri.fill.solid(); tri.fill.fore_color.rgb = PALEBLUE; _no_line(tri); tri.shadow.inherit = False
    tri2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, SW - Inches(2.6), SH - Inches(2.0),
                              Inches(2.6), Inches(2.0))
    tri2.fill.solid(); tri2.fill.fore_color.rgb = LIGHTBLUE; _no_line(tri2); tri2.shadow.inherit = False
    # 左上角小色块
    rect(s, Inches(0.9), Inches(0.9), Inches(0.7), Inches(0.14), BLUE)
    # 主标题
    tb, tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11), Inches(1.6))
    set_para(tf.paragraphs[0], "校园自习室预约系统", 46, True, DARKBLUE)
    # 副标题
    tb, tf = textbox(s, Inches(0.9), Inches(3.7), Inches(10.5), Inches(0.9))
    set_para(tf.paragraphs[0], "基于 Spring Cloud Alibaba 微服务架构与 AI 智能推荐",
             20, False, BLUE)
    # 分隔线
    rect(s, Inches(0.92), Inches(4.55), Inches(3.2), Pt(2.5), MIDBLUE)
    # 信息块
    info = [
        "课程：软件架构技术",
        "团队：CampusStudio",
        "成员：架构师/后端·AI  +  前端/文档·答辩",
        "指导教师：____      班级：____      日期：2026.06",
    ]
    tb, tf = textbox(s, Inches(0.92), Inches(4.85), Inches(9), Inches(2))
    for i, line in enumerate(info):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_para(p, line, 14, i == 0, INK if i else DARKBLUE, space_after=7)


# ---- P2 目录 ----
def slide_toc():
    s = new_slide()
    title(s, "目录", kicker="CONTENTS")
    items = [
        "项目背景与问题分析", "需求分析与功能设计",
        "系统架构设计（4+1 视图）", "数据库与国产化适配",
        "AI 智能模块设计", "云原生部署与验证",
        "AI 辅助开发实践", "项目总结与展望",
    ]
    # 两列
    col_w = Inches(5.6)
    for col in range(2):
        x = Inches(0.9) + col * Inches(6.2)
        tb, tf = textbox(s, x, Inches(2.05), col_w, Inches(4.3))
        for k in range(4):
            idx = col * 4 + k
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.space_after = Pt(20)
            r = p.add_run(); r.text = f"0{idx+1}   "
            _set_font(r, 22, True, BLUE)
            r2 = p.add_run(); r2.text = items[idx]
            _set_font(r2, 17, False, INK)


# ---- P3 背景 ----
def slide_bg():
    s = new_slide()
    title(s, "高校自习室资源管理的现实痛点", kicker="01  项目背景")
    cards(s, [
        ("！", "一座难求", "高峰期学生逐层找座，效率低下，体验差"),
        ("！", "占座不学", "空占现象普遍，资源严重浪费"),
        ("！", "考勤靠人工", "人工统计易出错、不及时"),
        ("！", "缺智能化", "无数据驱动的推荐与精细化管理"),
    ], Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.4))
    tb, tf = textbox(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.6))
    set_para(tf.paragraphs[0],
             "→ 亟需一套数字化、智能化、可精细管理的自习室预约系统",
             15, True, DARKBLUE, PP_ALIGN.CENTER)


# ---- P4 目标 ----
def slide_goal():
    s = new_slide()
    title(s, "五大核心目标", kicker="01  项目目标")
    items = [
        ("1", "资源数字化", "自习室、座位、时段全数字化建模"),
        ("2", "预约智能化", "AI 推荐最合适的座位与时段"),
        ("3", "考勤自动化", "签到签退自动记录与时长统计"),
        ("4", "管理精细化", "违规行为智能检测与申诉处理"),
        ("5", "架构云原生", "Docker + Kubernetes 容器化部署"),
    ]
    cards(s, items, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.6))


# ---- P5 需求总览 ----
def slide_req():
    s = new_slide()
    title(s, "功能模块与角色权限", kicker="02  需求分析")
    bullets(s, [
        ("7 大业务模块：用户管理 / 自习室管理 / 预约管理 / 考勤管理 / 违规管理 / AI 智能 / 通知管理", 0),
        ("3 类角色：学生、管理员、超级管理员（RBAC 分级授权）", 0),
        ("8 个核心用例：创建预约、签到、AI 推荐、取消预约、违规申诉、自习室管理、AI 异常分析、RAG 问答", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.6), Inches(4.6), size=14, gap=14)
    add_image(s, "diagram-05.png", Inches(6.7), Inches(1.95), Inches(5.9), Inches(4.5),
              caption="场景视图 · 用例图")


# ---- P6 技术选型 ----
def slide_tech():
    s = new_slide()
    title(s, "技术栈总览", kicker="03  技术选型")
    table(s, ["技术域", "技术选型", "版本"], [
        ["后端框架", "Spring Boot + Spring Cloud Alibaba", "JDK17 / Boot 3.2.5 / SCA 2023.0.1.2"],
        ["前端框架", "Vue3 + TypeScript + Pinia + Element Plus", "Vue 3.4.21"],
        ["数据库", "MySQL 8.0 + 达梦8（双库兼容）", "8.0.36 / DM8"],
        ["缓存", "Redis", "7.0.15"],
        ["服务治理", "Nacos 注册中心 + LoadBalancer", "2.3.0"],
        ["网关", "Spring Cloud Gateway", "动态路由 + JWT"],
        ["容器化", "Docker + Kubernetes", "Compose / K8s 1.29"],
    ], Inches(0.9), Inches(2.05), Inches(11.5), Inches(4.3),
       col_widths=[2, 5, 4.5])


# ---- P7 架构论证 ----
def slide_arch_choice():
    s = new_slide()
    title(s, "为什么选择微服务 + Spring Cloud Alibaba？", kicker="03  架构论证")
    bullets(s, [
        ("微服务 vs 单体：团队并行开发、独立部署、按服务弹性扩缩容", 0),
        ("本系统含 7 个清晰业务领域，天然适合微服务拆分", 1),
        ("SCA vs Netflix：Nacos / Sentinel / Seata 国产化生态完善", 0),
        ("Netflix OSS 已进入维护模式，SCA 更契合国内信创要求", 1),
        ("结论：采用 Spring Cloud Alibaba 微服务架构，兼顾扩展性与本土化", 0),
    ], Inches(0.9), Inches(2.1), Inches(11.4), Inches(4.4), size=16, gap=16)


# ---- P8 微服务拆分 ----
def slide_services():
    s = new_slide()
    title(s, "7 个微服务 + 1 网关 + 1 前端", kicker="03  服务拆分")
    table(s, ["服务", "端口", "核心职责"], [
        ["campus-gateway", "8000", "统一入口、JWT 校验、动态路由"],
        ["campus-auth", "8001", "登录注册、JWT 签发与刷新"],
        ["campus-user", "8002", "用户资料、角色权限"],
        ["campus-room", "8003", "自习室、座位、教学楼管理"],
        ["campus-reservation", "8004", "预约创建、取消、冲突检测"],
        ["campus-attendance", "8005", "考勤记录、学习时长统计"],
        ["campus-ai", "8006", "智能推荐、RAG 智能客服"],
    ], Inches(0.9), Inches(2.05), Inches(6.0), Inches(4.3),
       col_widths=[3.2, 1.2, 4.5], bsize=11)
    add_image(s, "diagram-04.png", Inches(7.1), Inches(2.0), Inches(5.5), Inches(4.3),
              caption="总体架构图")


# ---- P9 场景视图 ----
def slide_view_scenario():
    s = new_slide()
    title(s, "4+1 视图 · 场景视图：学生预约主线", kicker="04  架构设计")
    bullets(s, [
        ("访问前端 → 网关 JWT 认证", 0),
        ("浏览自习室列表、查看座位可用情况", 0),
        ("选择日期/时段/座位并提交预约", 0),
        ("预约服务执行时间冲突检测", 0),
        ("创建预约记录、更新座位状态", 0),
        ("预约时段内签到 / 签退", 0),
    ], Inches(0.9), Inches(2.05), Inches(4.7), Inches(4.5), size=14, gap=11)
    add_image(s, "diagram-09.png", Inches(5.8), Inches(1.95), Inches(6.8), Inches(4.6),
              caption="预约全链路时序图")


# ---- P10 逻辑视图 ----
def slide_view_logic():
    s = new_slide()
    title(s, "4+1 视图 · 逻辑视图：DDD 领域划分", kicker="04  架构设计")
    bullets(s, [
        ("5 大领域：用户域 / 空间域 / 预约域 / 考勤域 / AI 域", 0),
        ("核心实体：User、StudyRoom、StudySeat、Reservation、Attendance、KnowledgeBase", 0),
        ("设计模式：工厂（通知）、策略（推荐）、观察者（预约事件）", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.5), size=14, gap=13)
    add_image(s, "diagram-06.png", Inches(6.3), Inches(1.95), Inches(6.3), Inches(4.6),
              caption="领域模型类图")


# ---- P11 开发+进程视图 ----
def slide_view_dev():
    s = new_slide()
    title(s, "4+1 视图 · 开发视图 + 进程视图", kicker="04  架构设计")
    bullets(s, [
        ("Maven 多模块：backend / frontend / docs / k8s", 0),
        ("后端分层：Controller → Service → Mapper → Entity", 0),
        ("前端结构：api / views / components / stores / router", 0),
        ("运行时：前端 → Gateway → 微服务 → MySQL / Redis", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.0), Inches(4.5), size=14, gap=13)
    add_image(s, "diagram-22.png", Inches(6.1), Inches(1.95), Inches(6.5), Inches(4.6),
              caption="全链路调用时序图")


# ---- P12 物理视图 ----
def slide_view_physical():
    s = new_slide()
    title(s, "4+1 视图 · 物理视图：部署拓扑", kicker="04  架构设计")
    bullets(s, [
        ("11 个容器运行于 campus-network 桥接网络", 0),
        ("前端 nginx → Gateway → 7 微服务 → MySQL/Redis", 0),
        ("Nacos 负责服务注册与发现", 0),
        ("端口映射：MySQL 3307、Redis 6380、Nacos 8758", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.0), Inches(4.5), size=14, gap=13)
    add_image(s, "diagram-21.png", Inches(6.1), Inches(1.95), Inches(6.5), Inches(4.6),
              caption="Docker Compose 部署拓扑")


# ---- P13 安全架构 ----
def slide_security():
    s = new_slide()
    title(s, "认证、授权与数据安全", kicker="05  安全架构")
    bullets(s, [
        ("JWT + Refresh Token：accessToken 2 小时 / refreshToken 7 天", 0),
        ("RBAC 三级权限：学生 / 管理员 / 超级管理员", 0),
        ("密码 BCrypt 加密存储，传输层 HTTPS", 0),
        ("参数化 SQL 防注入 + XSS 过滤", 0),
        ("网关统一鉴权，Token 过期返回 401", 0),
    ], Inches(0.9), Inches(2.1), Inches(11.4), Inches(4.4), size=16, gap=15)


# ---- P14 数据库 ----
def slide_db():
    s = new_slide()
    title(s, "18 张表 + 双数据库兼容", kicker="05  数据库设计")
    bullets(s, [
        ("核心表：user / role / study_room / study_seat / reservation / attendance / violation / knowledge_base …", 0),
        ("遵循 3NF，统一审计字段（create_time / update_time / deleted）", 0),
        ("索引优化：reservation 复合索引 (room_id, seat_id, reserve_date)", 0),
        ("国产化：MySQL + 达梦8，JSON→CLOB、ENUM→VARCHAR", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.3), Inches(4.5), size=13, gap=13)
    add_image(s, "diagram-11.png", Inches(6.4), Inches(1.95), Inches(6.2), Inches(4.6),
              caption="数据库 ER 图")


# ---- P15 AI 推荐 ----
def slide_ai_rec():
    s = new_slide()
    title(s, "协同过滤 + 内容推荐 混合算法", kicker="06  AI 智能模块")
    bullets(s, [
        ("数据层：从 reservation 表构建用户-房间交互矩阵", 0),
        ("算法层：计算用户余弦相似度，取 Top-N 相似用户", 0),
        ("排序层：CF 分数(0.6) + 内容过滤分数(0.4) 加权融合", 0),
        ("解释层：调用智谱 GLM 生成自然语言推荐理由", 0),
        ("冷启动：新用户回退内容推荐 + 热门兜底", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.4), Inches(4.5), size=13, gap=11)
    add_image(s, "diagram-18.png", Inches(6.5), Inches(1.95), Inches(6.1), Inches(4.6),
              caption="智能推荐流程图")


# ---- P16 RAG ----
def slide_rag():
    s = new_slide()
    title(s, "基于知识库的 RAG 智能客服", kicker="06  AI 智能模块")
    bullets(s, [
        ("知识库：knowledge_base 表存储规则、流程、FAQ", 0),
        ("检索：关键词 + 中文 2-gram，title/keywords/content 加权打分", 0),
        ("生成：检索结果注入 system prompt，调用智谱 GLM-4", 0),
        ("降级：无 Key / API 失败时切换本地关键词匹配", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.4), Inches(4.5), size=13, gap=12)
    add_image(s, "diagram-19.png", Inches(6.5), Inches(1.95), Inches(6.1), Inches(4.6),
              caption="RAG 检索增强生成架构")


# ---- P17 AI 辅助开发 ----
def slide_ai_dev():
    s = new_slide()
    title(s, "AI 如何加速本项目开发", kicker="07  AI 辅助开发实践")
    bullets(s, [
        ("7 条真实 AI 辅助开发记录，累计节省 20+ 小时", 0),
        ("典型案例：协同过滤算法、RAG 检索、智谱 API 封装", 1),
        ("Docker 健康检查、数据库索引优化、Vue3 组件开发", 1),
        ("收益：开发效率提升、代码更规范、拓展知识边界", 0),
        ("把关：边界处理、业务上下文、安全与性能仍需人工审查", 0),
    ], Inches(0.9), Inches(2.1), Inches(11.4), Inches(4.4), size=15, gap=14)


# ---- P18 部署架构 ----
def slide_deploy():
    s = new_slide()
    title(s, "Docker Compose 一键部署", kicker="08  云原生部署")
    bullets(s, [
        ("11 个容器：3 基础(MySQL/Redis/Nacos) + 7 业务 + 1 前端", 0),
        ("健康检查 + depends_on: service_healthy 保证启动顺序", 0),
        ("预编译 jar + 轻量 JRE17 镜像，构建仅需 1–2 分钟", 0),
        ("一键命令：mvn clean package → docker compose up -d --build", 0),
    ], Inches(0.9), Inches(2.0), Inches(5.4), Inches(4.5), size=13, gap=12)
    add_screenshot(s, os.path.join(BASE, "screenshots", "docker", "02-docker-desktop.png"),
                   Inches(6.6), Inches(2.0), Inches(6.0), Inches(4.0),
                   caption="Docker Desktop 11 个容器 Running")


# ---- P19 部署验证 ----
def slide_verify():
    s = new_slide()
    title(s, "全链路验证通过", kicker="08  部署验证")
    table(s, ["验证项", "结果"], [
        ["11 个容器状态", "✅ 全部 Up / Healthy"],
        ["7 个服务注册 Nacos", "✅ 实例全部 healthy"],
        ["网关登录链路", "✅ 200 + JWT Token"],
        ["前端页面访问", "✅ http://localhost  HTTP 200"],
        ["AI 健康检查", '✅ {"status":"UP"}'],
    ], Inches(0.9), Inches(2.05), Inches(5.8), Inches(3.4),
       col_widths=[3, 4], bsize=13)
    add_screenshot(s, os.path.join(BASE, "screenshots", "docker", "01-docker-compose-ps.png"),
                   Inches(7.0), Inches(2.05), Inches(5.6), Inches(4.0),
                   caption="docker compose ps：全部 Up / Healthy")


# ---- P20 踩坑 ----
def slide_pitfall():
    s = new_slide()
    title(s, "5 个真实问题与排坑记录", kicker="08  部署踩坑")
    table(s, ["问题", "根因", "解决方案"], [
        ["镜像拉取超时", "Docker Hub 国内受限", "配置 6 个国内镜像加速器"],
        ["Nacos 端口绑定失败", "8848 落在 Windows 保留段", "宿主端口改为 8758"],
        ["服务启动崩溃", "config.import 无 dataId", "关配置中心，仅保留服务发现"],
        ["Redis 连接失败", "Boot3.2 配置项变更", "改用 spring.data.redis.host"],
        ["网关路由 503", "缺 loadbalancer 依赖", "网关 pom 增加 loadbalancer"],
    ], Inches(0.9), Inches(2.05), Inches(11.5), Inches(4.0),
       col_widths=[3, 3.5, 4.5], bsize=12)


# ---- P21 团队分工 ----
def slide_team():
    s = new_slide()
    title(s, "CampusStudio 团队分工", kicker="09  团队")
    cards(s, [
        ("70%", "架构师 / 后端 / AI",
         "负责 4+1 视图、微服务、数据库、AI 算法、Docker/K8s 部署"),
        ("30%", "前端 / 文档 / 答辩",
         "负责 Vue3 前端、前后端联调、文档、PPT、演示视频、Postman 测试"),
    ], Inches(1.4), Inches(2.6), Inches(10.5), Inches(3.2))


# ---- P22 亮点 ----
def slide_highlight():
    s = new_slide()
    title(s, "项目四大亮点与创新", kicker="09  亮点")
    cards(s, [
        ("①", "微服务完整落地", "7 业务服务 + 网关 + Nacos + LoadBalancer 动态路由"),
        ("②", "AI 双引擎", "协同过滤推荐 + RAG 客服，真实接入智谱 GLM-4"),
        ("③", "国产化适配", "MySQL + 达梦8 双库兼容，满足信创要求"),
        ("④", "云原生部署", "Docker Compose 跑通 11 容器，全链路验证"),
    ], Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.6))


# ---- P23 总结展望 ----
def slide_summary():
    s = new_slide()
    title(s, "总结与展望", kicker="09  总结")
    tb, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(5.5), Inches(0.5))
    set_para(tf.paragraphs[0], "✓  已完成", 16, True, BLUE)
    bullets(s, [
        ("核心预约 / 考勤 / 违规流程", 0),
        ("AI 推荐与 RAG 智能客服", 0),
        ("Docker Compose 全链路部署", 0),
        ("完整文档体系 + 100 页报告", 0),
    ], Inches(0.9), Inches(2.55), Inches(5.5), Inches(3.8), size=14, gap=11)
    tb, tf = textbox(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.5))
    set_para(tf.paragraphs[0], "→  未来演进", 16, True, MIDBLUE)
    bullets(s, [
        ("RabbitMQ 实现预约事件异步解耦", 0),
        ("SkyWalking / Prometheus / Grafana 可观测", 0),
        ("部署到 K8s 集群 + HPA 自动扩缩容", 0),
        ("接入学校统一认证、门禁硬件联动", 0),
    ], Inches(6.8), Inches(2.55), Inches(5.8), Inches(3.8), size=14, gap=11)


# ---- P24 致谢 ----
def slide_thanks():
    global _page
    s = prs.slides.add_slide(BLANK); _page += 1
    rect(s, 0, 0, SW, SH, DARKBLUE)
    # 装饰
    rect(s, 0, Inches(2.9), SW, Inches(0.06), MIDBLUE)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, SH - Inches(2.4),
                             Inches(3.4), Inches(2.4))
    tri.fill.solid(); tri.fill.fore_color.rgb = BLUE; _no_line(tri); tri.shadow.inherit = False
    tb, tf = textbox(s, Inches(0), Inches(2.9), SW, Inches(1.4), MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], "感谢聆听  ·  敬请指正", 40, True, WHITE, PP_ALIGN.CENTER)
    tb, tf = textbox(s, Inches(0), Inches(4.4), SW, Inches(0.6), MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], "Q  &  A", 24, False, LIGHTBLUE, PP_ALIGN.CENTER)
    tb, tf = textbox(s, Inches(0), Inches(5.4), SW, Inches(0.5), MSO_ANCHOR.MIDDLE)
    set_para(tf.paragraphs[0], "校园自习室预约系统  ·  CampusStudio", 13, False,
             LIGHTBLUE, PP_ALIGN.CENTER)


# ---------------- 主流程 ----------------
def main():
    slide_cover()
    slide_toc()
    slide_bg()
    slide_goal()
    slide_req()
    slide_tech()
    slide_arch_choice()
    slide_services()
    slide_view_scenario()
    slide_view_logic()
    slide_view_dev()
    slide_view_physical()
    slide_security()
    slide_db()
    slide_ai_rec()
    slide_rag()
    slide_ai_dev()
    slide_deploy()
    slide_verify()
    slide_pitfall()
    slide_team()
    slide_highlight()
    slide_summary()
    slide_thanks()
    prs.save(OUT)
    print(f"已生成 {OUT}，共 {_page} 页")


if __name__ == "__main__":
    main()
